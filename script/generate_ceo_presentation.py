from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(20, 45, 35)
ACCENT_COLOR = RGBColor(34, 139, 84)
BODY_COLOR = RGBColor(50, 60, 55)
LIGHT_BG = RGBColor(245, 249, 246)


def add_top_band(slide):
    band = slide.shapes.add_shape(
        autoshape_type_id=1,
        left=Inches(0),
        top=Inches(0),
        width=Inches(13.33),
        height=Inches(0.35),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT_COLOR
    band.line.fill.background()


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(1.5))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = TITLE_COLOR

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.6), Inches(2))
    stf = subtitle_box.text_frame
    stf.clear()
    p2 = stf.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = BODY_COLOR


def add_bullet_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = TITLE_COLOR

    content_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.8), Inches(5.6))
    ctf = content_box.text_frame
    ctf.clear()
    ctf.word_wrap = True

    first = True
    for bullet in bullets:
        paragraph = ctf.paragraphs[0] if first else ctf.add_paragraph()
        first = False
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = BODY_COLOR
        paragraph.space_after = Pt(10)

    if notes:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.clear()
        notes_frame.text = notes


def add_page_slide(prs, page_title, contains, message):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    tp = tf.paragraphs[0]
    tp.text = page_title
    tp.font.bold = True
    tp.font.size = Pt(28)
    tp.font.color.rgb = TITLE_COLOR

    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.3), Inches(5.6))
    ltf = left_box.text_frame
    ltf.clear()
    head = ltf.paragraphs[0]
    head.text = "What This Page Contains"
    head.font.bold = True
    head.font.size = Pt(20)
    head.font.color.rgb = ACCENT_COLOR
    head.space_after = Pt(8)

    for item in contains:
        p = ltf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(6)

    right_shape = slide.shapes.add_shape(
        autoshape_type_id=1,
        left=Inches(9.3),
        top=Inches(1.6),
        width=Inches(3.5),
        height=Inches(3.0),
    )
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = RGBColor(235, 245, 238)
    right_shape.line.color.rgb = ACCENT_COLOR

    rtf = right_shape.text_frame
    rtf.clear()
    rp1 = rtf.paragraphs[0]
    rp1.text = "Key Message"
    rp1.font.bold = True
    rp1.font.size = Pt(18)
    rp1.font.color.rgb = TITLE_COLOR
    rp1.alignment = PP_ALIGN.CENTER

    rp2 = rtf.add_paragraph()
    rp2.text = message
    rp2.font.size = Pt(15)
    rp2.font.color.rgb = BODY_COLOR

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = f"Say: {message}"


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "CADI Website Executive Presentation",
        "CEO Walkthrough | Operations Transparency Platform\nDate: February 25, 2026",
    )

    add_bullet_slide(
        prs,
        "Agenda",
        [
            "1. Strategic architecture overview",
            "2. Page-by-page walkthrough",
            "3. Key executive takeaways",
            "4. Q&A",
        ],
    )

    add_bullet_slide(
        prs,
        "Architecture Framework",
        [
            "Credibility and Entry: Dashboard Overview, About CADI",
            "Operations and Delivery: Nursery, Plantation, Client Services, Updates Log, Photo Gallery",
            "Risk and Positioning: Zambales Location, Compliance and Regulations",
            "Institutional Capability: AI Technology, Community Impact, Management Team",
        ],
        notes="This structure mirrors board-level due diligence flow.",
    )

    pages = [
        (
            "Dashboard Overview",
            [
                "Hero section with value proposition and action buttons",
                "Live nursery metrics: seedling counts, heights, mortality, out-planting target",
                "Operational updates preview feed",
                "Featured operations gallery: nursery, land, out-planting readiness",
                "Zambales strategic highlight and core values cards",
            ],
            "This is the executive command center for strategy, activity, and trust.",
        ),
        (
            "About CADI",
            [
                "Corporate profile and management mandate",
                "Corporate overview and BGC positioning",
                "Operating mandate pillars: management, compliance, precision agriculture",
                "Governance structure visual: brand and sales versus operations execution",
            ],
            "This page clarifies accountability and governance.",
        ),
        (
            "Nursery Operations",
            [
                "Live seedling gallery with staged growth visuals",
                "Stock propagation details for Aquilaria and Sweet Elena Mango",
                "Growth dashboard with metrics and update timestamp",
                "Nursery technology protocol: irrigation, climate, pest, soil analytics",
            ],
            "This page demonstrates biological pipeline quality and control.",
        ),
        (
            "Plantation Operations",
            [
                "Leased land and preparation gallery",
                "July 2026 out-planting block with projected volumes",
                "Land preparation cards: soil, spacing geometry, cassava intercropping",
                "Lifecycle timeline from establishment to harvest",
            ],
            "This page shows field execution readiness and long-range planning.",
        ),
        (
            "Client Services",
            [
                "Individual tree tracking: serialized ID, GPS mapping, portal access, audits",
                "Professional reporting: annual reports, drone analytics, health, yield forecasts",
                "Client visitation program with logistics and two-day itinerary",
            ],
            "This page turns ownership into verifiable and reportable assets.",
        ),
        (
            "Updates Log",
            [
                "Full archive of operational updates",
                "Search and category filters",
                "Dated cards with tags and optional images",
                "Summary stats: total updates, categories, latest update",
            ],
            "This page proves operational cadence and continuity.",
        ),
        (
            "Photo Gallery",
            [
                "Consolidated photo library with search",
                "Category tabs: nursery, plantation, facilities, operations",
                "Lightbox with navigation and metadata",
                "Gallery summary statistics",
            ],
            "This page provides visual evidence of field activity and transparency.",
        ),
        (
            "Zambales Location",
            [
                "Geographic focus panel for Southern Zambales",
                "Logistical hub details: airport, roads, seaport",
                "Agro-climatic profile: dry season, rainfall, soil",
                "Strategic split: agricultural excellence and ecotourism gateway",
            ],
            "This page explains why location choice is strategically strong.",
        ),
        (
            "Compliance and Regulations",
            [
                "Executive compliance advantage panel",
                "Regulatory framework cards: DENR, CITES, Customs, Phytosanitary",
                "Further mandates: PEFC pathway and land lease security",
            ],
            "This page addresses legal risk and export readiness.",
        ),
        (
            "AI Technology",
            [
                "Technology stack: drones, IoT sensors, smart irrigation, inoculation formulas",
                "Operational benefits: mortality reduction, cost optimization, early warning, forecasting",
            ],
            "This page positions CADI as a data-enabled operator.",
        ),
        (
            "Community Impact",
            [
                "1:1 reforestation commitment",
                "Local employment and skills training",
                "Cassava intercropping donation commitment",
            ],
            "This page links commercial execution with community value.",
        ),
        (
            "Management Team",
            [
                "Leadership overview",
                "Grouped team directory: executive management, board, senior management",
                "Profile details: role, experience, technical expertise",
            ],
            "This page closes the trust loop with leadership depth.",
        ),
    ]

    for title, contains, message in pages:
        add_page_slide(prs, title, contains, message)

    add_bullet_slide(
        prs,
        "Strategic Close",
        [
            "Clear governance and accountability",
            "Verifiable operational evidence",
            "Structured compliance and risk management",
            "Technology-backed execution and institutional leadership",
            "CADI's website is the digital operating face of the organization",
        ],
    )

    add_bullet_slide(
        prs,
        "Q&A",
        [
            "Operational readiness and deployment timeline",
            "Regulatory security and compliance model",
            "Reporting transparency and stakeholder visibility",
        ],
    )

    output_path = Path("CADI_Website_Executive_Presentation.pptx")
    prs.save(output_path)
    print(f"Created: {output_path.resolve()}")


if __name__ == "__main__":
    main()
