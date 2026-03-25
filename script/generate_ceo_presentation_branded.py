from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PRIMARY = RGBColor(25, 102, 68)     # CADI institutional green
ACCENT = RGBColor(198, 112, 6)      # Earthy accent
TEXT_DARK = RGBColor(34, 41, 38)
TEXT_MUTED = RGBColor(74, 85, 80)
BG_LIGHT = RGBColor(246, 249, 247)
WHITE = RGBColor(255, 255, 255)

LOGO_PATH = Path("client/public/golden-forests-logo.png")


def apply_brand_chrome(slide, title_slide=False):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_LIGHT

    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.32)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = PRIMARY
    top.line.fill.background()

    bottom = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.26), Inches(13.33), Inches(0.24)
    )
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = RGBColor(232, 240, 236)
    bottom.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.26), Inches(4.3), Inches(0.24)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(11.95), Inches(0.38), height=Inches(0.55))

    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.25), Inches(7), Inches(0.24))
    ft = footer.text_frame
    ft.clear()
    p = ft.paragraphs[0]
    p.text = "CADI Operations | Golden Forests Philippines"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE if title_slide else PRIMARY


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_brand_chrome(slide, title_slide=True)

    # subtle hero panel
    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.12), Inches(11.8), Inches(4.95)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = WHITE
    hero.line.color.rgb = RGBColor(214, 227, 220)

    tbox = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10.8), Inches(1.6))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "CADI Website Executive Presentation"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.LEFT

    sbox = slide.shapes.add_textbox(Inches(1.2), Inches(3.15), Inches(10.4), Inches(2.0))
    stf = sbox.text_frame
    stf.clear()
    p2 = stf.paragraphs[0]
    p2.text = "CEO Walkthrough\nOperations Transparency Platform\nFebruary 25, 2026"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_MUTED


def add_bullets_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_brand_chrome(slide)

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.62), Inches(11.5), Inches(0.8))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = PRIMARY

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.55), Inches(11.8), Inches(5.45)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(214, 227, 220)

    box = slide.shapes.add_textbox(Inches(1.15), Inches(1.95), Inches(11.0), Inches(4.9))
    btf = box.text_frame
    btf.clear()
    btf.word_wrap = True

    first = True
    for b in bullets:
        para = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        para.text = b
        para.font.size = Pt(22)
        para.font.color.rgb = TEXT_DARK
        para.space_after = Pt(11)

    if notes:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.clear()
        notes_tf.text = notes


def add_page_slide(prs, title, contains, key_message):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_brand_chrome(slide)

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.58), Inches(11.6), Inches(0.8))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = PRIMARY

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(8.35), Inches(5.5)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = RGBColor(214, 227, 220)

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(1.5), Inches(3.2), Inches(5.5)
    )
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(236, 244, 239)
    right.line.color.rgb = PRIMARY

    hbox = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(7.6), Inches(0.5))
    htf = hbox.text_frame
    htf.clear()
    h = htf.paragraphs[0]
    h.text = "What This Page Contains"
    h.font.bold = True
    h.font.size = Pt(20)
    h.font.color.rgb = ACCENT

    cbox = slide.shapes.add_textbox(Inches(1.1), Inches(2.25), Inches(7.7), Inches(4.5))
    ctf = cbox.text_frame
    ctf.clear()
    ctf.word_wrap = True

    first = True
    for item in contains:
        cp = ctf.paragraphs[0] if first else ctf.add_paragraph()
        first = False
        cp.text = item
        cp.font.size = Pt(17)
        cp.font.color.rgb = TEXT_DARK
        cp.space_after = Pt(6)

    kbox = slide.shapes.add_textbox(Inches(9.7), Inches(1.95), Inches(2.6), Inches(0.6))
    ktf = kbox.text_frame
    ktf.clear()
    kp = ktf.paragraphs[0]
    kp.text = "Key Message"
    kp.font.bold = True
    kp.font.size = Pt(18)
    kp.font.color.rgb = PRIMARY
    kp.alignment = PP_ALIGN.CENTER

    mb = slide.shapes.add_textbox(Inches(9.7), Inches(2.65), Inches(2.6), Inches(3.95))
    mtf = mb.text_frame
    mtf.clear()
    mp = mtf.paragraphs[0]
    mp.text = key_message
    mp.font.size = Pt(16)
    mp.font.color.rgb = TEXT_DARK
    mp.alignment = PP_ALIGN.LEFT

    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = f"Say: {key_message}"


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullets_slide(
        prs,
        "Agenda",
        [
            "1. Strategic architecture overview",
            "2. Page-by-page website walkthrough",
            "3. Executive takeaways",
            "4. Q&A",
        ],
    )

    add_bullets_slide(
        prs,
        "Architecture Framework",
        [
            "Credibility and Entry: Dashboard Overview, About CADI",
            "Operations and Delivery: Nursery Operations, Plantation Operations, Client Services, Updates Log, Photo Gallery",
            "Risk and Positioning: Zambales Location, Compliance and Regulations",
            "Institutional Capability: AI Technology, Community Impact, Management Team",
        ],
        notes="Use this slide to set the business logic before opening individual pages.",
    )

    page_slides = [
        ("Dashboard Overview", [
            "Hero section with value proposition and primary action buttons",
            "Live nursery metrics: seedling counts, heights, mortality, out-planting target",
            "Latest operational updates preview feed",
            "Featured operations gallery: nursery, land, out-planting readiness",
            "Zambales strategic highlight and core values cards",
        ], "Executive command center for strategy, activity, and trust."),
        ("About CADI", [
            "Corporate profile header and management mandate",
            "Corporate overview and BGC positioning",
            "Operating mandate pillars: management, compliance, precision agriculture",
            "Governance visual: brand and sales versus operations execution",
        ], "Clarifies accountability and governance structure."),
        ("Nursery Operations", [
            "Live seedling gallery with staged growth visuals",
            "Stock propagation details for Aquilaria and Sweet Elena Mango",
            "Growth dashboard with metrics and update timestamp",
            "Technology protocol: irrigation, climate control, pest management, soil analytics",
        ], "Demonstrates biological pipeline quality and control."),
        ("Plantation Operations", [
            "Leased land and preparation gallery",
            "July 2026 out-planting block with projected deployment volumes",
            "Land preparation: soil protocol, spacing geometry, cassava intercropping",
            "Lifecycle timeline from establishment to harvest",
        ], "Shows field execution readiness and lifecycle planning."),
        ("Client Services", [
            "Tree tracking: serialized IDs, GPS mapping, portal access, audit verification",
            "Professional reporting: annual reports, drone analytics, health metrics, yield forecasts",
            "Client visitation programme with logistics and itinerary",
        ], "Turns ownership into verifiable and reportable assets."),
        ("Updates Log", [
            "Full archive of operational updates",
            "Search and category filter controls",
            "Dated update cards with category tags and optional images",
            "Summary statistics: total updates, categories, latest date",
        ], "Proves operational cadence and continuity."),
        ("Photo Gallery", [
            "Consolidated operational gallery with search",
            "Category tabs: nursery, plantation, facilities, operations",
            "Lightbox with navigation and image metadata",
            "Gallery summary stats by category",
        ], "Provides visual evidence of real field activity."),
        ("Zambales Location", [
            "Geographic focus panel for Southern Zambales",
            "Logistics hub details: airport, roads, seaport proximity",
            "Agro-climatic profile: seasonality, rainfall, soil",
            "Strategic value: agricultural excellence and ecotourism gateway",
        ], "Explains strategic strength of site selection."),
        ("Compliance and Regulations", [
            "Executive compliance advantage panel",
            "Core framework: DENR, CITES, Customs, Phytosanitary",
            "Further mandates: PEFC pathway and long-term lease security",
        ], "Addresses legal risk and export readiness."),
        ("AI Technology", [
            "Technology stack: drones, IoT sensors, smart irrigation, inoculation formulas",
            "Operational benefits: lower mortality, optimized inputs, early warnings, stronger forecasting",
        ], "Positions CADI as a modern data-enabled operator."),
        ("Community Impact", [
            "1:1 reforestation programme",
            "Local employment and skills training",
            "Cassava intercropping donation commitment",
        ], "Links commercial performance with social and environmental value."),
        ("Management Team", [
            "Leadership overview",
            "Grouped team directory: executive management, board, senior management",
            "Profiles with role, experience, and expertise",
        ], "Closes trust loop with leadership depth."),
    ]

    for title, contains, message in page_slides:
        add_page_slide(prs, title, contains, message)

    add_bullets_slide(
        prs,
        "Strategic Close",
        [
            "Clear governance and accountability",
            "Verifiable operational evidence",
            "Structured compliance and risk management",
            "Technology-backed execution and institutional leadership",
            "CADI website functions as the digital operating face of the organization",
        ],
    )

    add_bullets_slide(
        prs,
        "Q&A",
        [
            "Operational readiness and deployment timeline",
            "Regulatory security and compliance model",
            "Reporting transparency and stakeholder visibility",
        ],
    )

    out = Path("CADI_Website_Executive_Presentation_BRANDED.pptx")
    prs.save(out)
    print(f"Created: {out.resolve()}")


if __name__ == "__main__":
    main()
